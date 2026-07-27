"""
PPR pipeline - Stage 5: one PowerPoint per center, ready for the review meeting.

Kolin's bottleneck (Meet 11): "I have about 6 to 8 of these I need to produce in the
next two days. And I don't have a gap in my calendar." The dashboard makes the numbers
easy to read; this makes the decks he actually presents.

Each deck mirrors the real per-center reviews:
    Slide 1  Launch-to-Date Metrics   center vs Top 10 / Top 40 / national median
    Slide 2  Year over Year Metrics   two periods side by side with the difference

Usage:
    python build_center_decks.py                 all centers
    python build_center_decks.py "Moffitt" "Yale"  named centers (substring match)

Out: decks/<Center Name> - P&PR Review.pptx
"""
import os
import sys

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
ANA = os.path.join(HERE, "..", "analysis")
OUT_DIR = os.path.join(HERE, "..", "decks")

NAVY = RGBColor(0x17, 0x34, 0x4F)
LIME = RGBColor(0x9D, 0xC1, 0x3C)
INK = RGBColor(0x26, 0x33, 0x3D)
MUTE = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xEA, 0xF0, 0xE4)
LINE = RGBColor(0xD0, 0xD9, 0xE0)
GOOD = RGBColor(0x2E, 0x7D, 0x32)
BAD = RGBColor(0xC0, 0x39, 0x2B)
FONT = "Arial"

# Lower is better, so a negative difference is an improvement.
LOWER_IS_BETTER = {
    "TTPs Cancelled or Rescheduled within 7 Days Prior to Slot Reservation",
    "Patient Related Drop-outs following TTP due to patient health",
    "OOS Products", "Patient Progression Rate",
    "Average Time From Enrollment Date to TTP (Days)",
    "Average Time From TTP to AMTAGVI Infusion (Days)",
    "Average Time From Final Product Delivery Date to AMTAGVI Infusion (Days)",
}


def fmt(value, value_type):
    if pd.isna(value):
        return "-"
    if value_type == "rate":
        return f"{value * 100:.1f}%"
    if value_type == "days":
        return f"{value:.1f}"
    return f"{int(round(value)):,}"


def fmt_delta(value, value_type):
    if pd.isna(value):
        return "-"
    if abs(value) < 1e-9:
        return "-"
    sign = "+" if value > 0 else ""
    if value_type == "rate":
        return f"{sign}{value * 100:.1f}pt"
    if value_type == "days":
        return f"{sign}{value:.1f}"
    return f"{sign}{int(round(value))}"


def add_textbox(slide, x, y, w, h, text, size, *, bold=False, color=INK,
                align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    return box


def add_band(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def chrome(slide, eyebrow, title, asof):
    """Iovance frame: eyebrow, title, footer band, confidentiality line."""
    add_textbox(slide, 0.55, 0.34, 11, 0.25, eyebrow.upper(), 10.5, bold=True, color=NAVY)
    add_textbox(slide, 0.55, 0.62, 11.5, 0.55, title, 26, bold=True, color=NAVY)
    add_band(slide, 0, 6.95, 13.333, 0.55, LIME)
    add_textbox(slide, 0.55, 7.08, 8, 0.3,
                "© 2026, Iovance Biotherapeutics, Inc.  |  Confidential for Internal Use Only",
                8, color=NAVY)
    add_textbox(slide, 9.2, 7.08, 3.6, 0.3, f"Data as of {asof}", 8,
                color=NAVY, align=PP_ALIGN.RIGHT)


def style_table(table, col_widths, header_rows=1):
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = Emu(int(Inches(width)))
    for r, row in enumerate(table.rows):
        row.height = Emu(int(Inches(0.32)))
        for c, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            if r < header_rows:
                cell.fill.fore_color.rgb = NAVY
            elif r % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
            else:
                cell.fill.fore_color.rgb = WHITE


def write_cell(cell, text, *, size=10, bold=False, color=INK, align=PP_ALIGN.RIGHT):
    frame = cell.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def slide_launch_to_date(prs, center, sc, asof):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(slide, f"{center} - Patient & Process Review", "Launch-to-Date Metrics", asof)

    cols = ["Launch to Date", "Top 10", "Top 40", "New"]
    headers = ["Category", "Metric", "Launch to Date", "Top 10 ATCs", "Top 40 ATCs", "'New' ATCs"]
    metrics = sc.sort_values("metric_order").metric.unique().tolist()

    shape = slide.shapes.add_table(len(metrics) + 1, len(headers),
                                   Inches(0.55), Inches(1.45),
                                   Inches(12.2), Inches(0.32 * (len(metrics) + 1)))
    table = shape.table
    style_table(table, [1.9, 4.1, 1.6, 1.55, 1.55, 1.5])

    for c, head in enumerate(headers):
        write_cell(table.cell(0, c), head, size=10, bold=True, color=WHITE,
                   align=PP_ALIGN.LEFT if c < 2 else PP_ALIGN.RIGHT)

    prev_group = None
    for r, metric in enumerate(metrics, start=1):
        rows = sc[sc.metric == metric]
        group = rows.metric_group.iloc[0]
        vtype = rows.value_type.iloc[0]
        write_cell(table.cell(r, 0), "" if group == prev_group else group,
                   size=8.5, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
        prev_group = group
        write_cell(table.cell(r, 1), metric, size=9, align=PP_ALIGN.LEFT)
        for c, col in enumerate(cols, start=2):
            match = rows[rows.col_label == col]
            value = match.value.iloc[0] if len(match) else float("nan")
            write_cell(table.cell(r, c), fmt(value, vtype), size=10,
                       bold=(c == 2), color=NAVY if c == 2 else MUTE)

    add_textbox(slide, 0.55, 6.35, 12, 0.5,
                "Patient Progression Rate = (patient related drop-offs after mfg. start) / (mfg. starts). "
                "Benchmarks are the median across centers in each tier; timing metrics are medians.",
                8, italic=True, color=MUTE)


def slide_year_over_year(prs, center, sc, asof):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(slide, f"{center} - Patient & Process Review", "Year over Year Metrics at ATC", asof)

    cols = ["2025", "2026 YTD"]
    headers = ["Category", "Metric", "2025", "2026 YTD", "Difference"]
    metrics = sc.sort_values("metric_order").metric.unique().tolist()

    shape = slide.shapes.add_table(len(metrics) + 1, len(headers),
                                   Inches(0.55), Inches(1.45),
                                   Inches(12.2), Inches(0.32 * (len(metrics) + 1)))
    table = shape.table
    style_table(table, [1.9, 4.7, 1.9, 1.9, 1.8])

    for c, head in enumerate(headers):
        write_cell(table.cell(0, c), head, size=10, bold=True, color=WHITE,
                   align=PP_ALIGN.LEFT if c < 2 else PP_ALIGN.RIGHT)

    prev_group = None
    for r, metric in enumerate(metrics, start=1):
        rows = sc[sc.metric == metric]
        group = rows.metric_group.iloc[0]
        vtype = rows.value_type.iloc[0]
        write_cell(table.cell(r, 0), "" if group == prev_group else group,
                   size=8.5, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
        prev_group = group
        write_cell(table.cell(r, 1), metric, size=9, align=PP_ALIGN.LEFT)

        values = []
        for c, col in enumerate(cols, start=2):
            match = rows[rows.col_label == col]
            value = match.value.iloc[0] if len(match) else float("nan")
            values.append(value)
            write_cell(table.cell(r, c), fmt(value, vtype), size=10, color=INK)

        delta = values[1] - values[0] if not any(pd.isna(v) for v in values) else float("nan")
        color = MUTE
        if not pd.isna(delta) and abs(delta) > 1e-9:
            improved = delta < 0 if metric in LOWER_IS_BETTER else delta > 0
            color = GOOD if improved else BAD
        write_cell(table.cell(r, 4), fmt_delta(delta, vtype), size=10, bold=True, color=color)

    add_textbox(slide, 0.55, 6.35, 12, 0.5,
                "Each metric is counted on its own event date (enrollment, TTP, delivery, or infusion), "
                "so a column reflects what happened in that period. Green is an improvement.",
                8, italic=True, color=MUTE)


def build_deck(center, sc_center, national, asof):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    combined = pd.concat([sc_center, national], ignore_index=True)
    slide_launch_to_date(prs, center, combined, asof)
    slide_year_over_year(prs, center, sc_center, asof)
    safe = "".join(ch for ch in center if ch.isalnum() or ch in " -_&.").strip()
    path = os.path.join(OUT_DIR, f"{safe} - P&PR Review.pptx")
    prs.save(path)
    return path


def main():
    tidy = pd.read_csv(os.path.join(ANA, "ppr_scorecard_tidy.csv"))
    asof = os.environ.get("PPR_ASOF", "2026-07-21")
    os.makedirs(OUT_DIR, exist_ok=True)

    national = tidy[tidy.scope == "National"]
    centers = sorted(tidy[tidy.scope == "Center"].center.unique())
    if len(sys.argv) > 1:
        wanted = [a.lower() for a in sys.argv[1:]]
        centers = [c for c in centers if any(w in c.lower() for w in wanted)]
        if not centers:
            print("No center matched:", ", ".join(sys.argv[1:]))
            return 1

    for center in centers:
        build_deck(center, tidy[(tidy.scope == "Center") & (tidy.center == center)],
                   national, asof)
    print(f"{len(centers)} deck(s) -> {os.path.abspath(OUT_DIR)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
