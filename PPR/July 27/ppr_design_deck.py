"""
P&PR scorecard: table design proposal deck.

Built to the Iovance house style reverse-engineered from the real deck
"2H'26 AMTAGVI CTAM_RAD IC Overviews" and from Kolin's own (Proposed) P&PR
Metrics.xlsx template photographed in PPR Automation/Goal/IMG_9401.

Design rules taken from those two sources, not invented:
  - content slide: steel-blue eyebrow, navy full-sentence action title,
    two small olive squares flanking the title, olive footer band
  - table: grey block header row, light-green time-column headers,
    national block boxed in heavy black, category column shaded in blues,
    thin black grid on every cell, no row banding
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

C = dict(
    navy=RGBColor(0x17, 0x34, 0x4F),
    steel=RGBColor(0x2F, 0x5D, 0x8A),
    lime=RGBColor(0x9D, 0xC1, 0x3C),
    olive=RGBColor(0x56, 0x7A, 0x2E),
    red=RGBColor(0xC0, 0x39, 0x2B),
    white=RGBColor(0xFF, 0xFF, 0xFF),
    black=RGBColor(0x00, 0x00, 0x00),
    ink=RGBColor(0x26, 0x33, 0x3D),
    grey_hdr=RGBColor(0xD9, 0xD9, 0xD9),
    green_hdr=RGBColor(0xD8, 0xE4, 0xBC),
    blue1=RGBColor(0xDC, 0xE6, 0xF1),
    blue2=RGBColor(0xC5, 0xD9, 0xF1),
    blue3=RGBColor(0xB8, 0xCC, 0xE4),
    faint=RGBColor(0x8A, 0x97, 0xA1),
    panel=RGBColor(0xF2, 0xF5, 0xF8),
)
FONT = "Segoe UI"
W, H = Inches(13.333), Inches(7.5)


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(s, x, y, w, h, fill=None, line=None, lw=0.75):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0):
    """runs: list of (string, size, bold, color, italic) or a list of such lists
    (one inner list per paragraph)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        pg = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pg.alignment = align
        if space:
            pg.space_after = Pt(space)
        for item in para:
            t, sz, bold, col = item[0], item[1], item[2], item[3]
            it = item[4] if len(item) > 4 else False
            r = pg.add_run(); r.text = t
            r.font.name = FONT; r.font.size = Pt(sz); r.font.bold = bold
            r.font.italic = it; r.font.color.rgb = col
    return tb


def chrome(s, eyebrow, title, page):
    """The Iovance content-slide frame: eyebrow, action title, two olive
    squares flanking it, olive footer band with the legal line and wordmark."""
    text(s, Inches(0.62), Inches(0.42), Inches(11.5), Inches(0.3),
         [(eyebrow, 11, True, C["steel"])])
    text(s, Inches(0.62), Inches(0.74), Inches(11.55), Inches(0.86),
         [(title, 21, True, C["navy"])])
    for xx in (Inches(0.28), Inches(12.78)):
        box(s, xx, Inches(0.80), Inches(0.19), Inches(0.19), fill=C["olive"])
    box(s, 0, H - Inches(0.40), W, Inches(0.40), fill=C["olive"])
    text(s, Inches(0.30), H - Inches(0.32), Inches(7), Inches(0.25),
         [("© 2025, Iovance Biotherapeutics, Inc.  |  Confidential for Internal Use Only",
           8.5, False, C["white"])])
    text(s, Inches(10.9), H - Inches(0.33), Inches(1.7), Inches(0.26),
         [("I O V A N C E", 11, True, C["white"])], align=PP_ALIGN.RIGHT)
    box(s, Inches(12.78), H - Inches(0.33), Inches(0.26), Inches(0.26),
        fill=C["white"])
    text(s, Inches(12.78), H - Inches(0.30), Inches(0.26), Inches(0.2),
         [(str(page), 8, True, C["olive"])], align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------- slide 1
def s1(prs):
    s = blank(prs)
    box(s, 0, 0, W, H, fill=C["navy"])
    box(s, Inches(8.6), 0, Inches(4.73), Inches(1.15), fill=C["lime"])
    box(s, 0, H - Inches(0.92), W, Inches(0.62), fill=C["lime"])
    text(s, Inches(0.9), Inches(1.5), Inches(7), Inches(0.9),
         [[("I", 40, True, C["white"]), ("O", 40, True, C["lime"]),
           ("VANCE", 40, True, C["white"])],
          [("B I O T H E R A P E U T I C S", 12, True, C["lime"])]])
    text(s, Inches(0.9), Inches(3.3), Inches(10.6), Inches(1.8),
         [("Redesigning the P&PR scorecard so a centre review reads in one look",
           30, True, C["white"])])
    text(s, Inches(0.9), Inches(5.1), Inches(8), Inches(0.4),
         [("Business Analytics and Insights  |  July 2026", 12, False, C["lime"])])
    text(s, Inches(0.9), H - Inches(0.80), Inches(9), Inches(0.4),
         [("A D V A N C I N G   I M M U N O - O N C O L O G Y", 13, True, C["white"])])
    text(s, Inches(0.9), H - Inches(0.26), Inches(9), Inches(0.2),
         [("© 2025, Iovance Biotherapeutics, Inc.  |  Confidential for Internal Use Only",
           8, False, C["white"])])


# ---------------------------------------------------------------- slide 2
def s2(prs):
    s = blank(prs)
    chrome(s, "P&PR Scorecard  |  The problem today",
           "One centre review takes over an hour by hand, and two reps can get two different answers", 2)
    steps = [("Download", "Pull the balanced scorecard\nand six more reports out of\nInfinity, one at a time"),
             ("Paste", "Drop them into an Excel that\ncomputes launch-to-date\nagainst quartiles"),
             ("Rebuild", "Retype the numbers onto a\nslide, reshaped for whichever\nregion asked"),
             ("Repeat", "Do it again for the next\ncentre. Eighty-five centres,\nno two decks alike")]
    x = Inches(0.62)
    for i, (h, b) in enumerate(steps):
        box(s, x, Inches(2.15), Inches(2.85), Inches(1.42),
            fill=C["panel"], line=C["faint"], lw=0.5)
        text(s, x + Inches(0.22), Inches(2.32), Inches(2.4), Inches(0.3),
             [(h, 14, True, C["olive"])])
        text(s, x + Inches(0.22), Inches(2.72), Inches(2.5), Inches(0.8),
             [(b, 11, False, C["ink"])])
        x += Inches(3.02)
    box(s, Inches(0.62), Inches(4.10), Inches(11.9), Inches(1.42),
        fill=C["white"], line=C["navy"], lw=1.25)
    text(s, Inches(0.95), Inches(4.30), Inches(11.3), Inches(1.1),
         [[("What it costs", 13, True, C["navy"])],
          [("About an hour and a half per centre, and the output is not uniform. "
            "In Kolin's words the reps get different answers, he gives one number and "
            "somebody else gives another. The quartile columns were meant to help and "
            "they do the opposite: he says they confuse the sales folks and they confuse "
            "the people in the treatment centres, and the team is actively moving away "
            "from them.", 12, False, C["ink"])]], space=6)


# ---------------------------------------------------------------- slide 3
def s3(prs):
    s = blank(prs)
    chrome(s, "P&PR Scorecard  |  Who reads it",
           "A centre asks three questions, and the table has to answer all three on its own", 3)
    qs = [("How am I doing?", "Their own volume and timings,\nlaunch to date and by year",
           "Launch to Date, 2024, 2025,\n2026 YTD"),
          ("Compared with whom?", "One blinded peer group, matched\nto their size. Never a named list\nof other centres",
           "Top 10, Top 40, or New.\nOne arm, not all three"),
          ("Am I getting better?", "The last four quarters, so a centre\ncan see its own direction rather\nthan only its rank",
           "Q3'26 QTD, Q2'26,\nQ1'26, Q4'25")]
    x = Inches(0.62)
    for q, b, cols in qs:
        box(s, x, Inches(2.0), Inches(3.85), Inches(2.75),
            fill=C["white"], line=C["olive"], lw=1.0)
        text(s, x + Inches(0.28), Inches(2.25), Inches(3.3), Inches(0.4),
             [(q, 16, True, C["navy"])])
        text(s, x + Inches(0.28), Inches(2.78), Inches(3.3), Inches(1.0),
             [(b, 11.5, False, C["ink"])])
        box(s, x + Inches(0.28), Inches(3.82), Inches(3.3), Inches(0.72),
            fill=C["panel"])
        text(s, x + Inches(0.45), Inches(3.95), Inches(3.0), Inches(0.6),
             [(cols, 10.5, True, C["olive"])])
        x += Inches(4.03)
    text(s, Inches(0.62), Inches(5.15), Inches(11.9), Inches(1.0),
         [("The template already answers all three. What it does not do is survive being "
           "rebuilt by hand every time, which is where the differing answers come from.",
           13, False, C["ink"])])


# ---------------------------------------------------------------- the table
METRICS = [
    ("Patient Identification\n& Enrollment", "Enrollments in IovanceCares", "11", "", "8", "3", "36", "29", "20", "0", "2", "1", "5"),
    (None, "Patients Enrolled in IovanceCares", "11", "", "8", "3", "35", "28", "19", "0", "2", "1", "5"),
    (None, "TTPs Cancelled or Rescheduled\n<=7 Days Prior to Slot Reservation *", "2", "", "1", "1", "5", "4", "2", "0", "1", "0", "1"),
    ("Tumor Tissue\nProcurement", "Completed TTPs", "8", "", "3", "5", "21", "18", "10", "0", "1", "4", "2"),
    (None, "Scheduled TTPs", "0", "", "0", "0", "2", "1", "0", "0", "0", "0", "0"),
    (None, "2nd Resections (Scheduled or Completed)", "0", "", "0", "0", "2", "1", "0", "0", "0", "0", "0"),
    ("AMTAGVI\nRegimen", "Patient Related Drop-outs following\nTTP due to patient health", "1", "", "", "1", "2", "1", "0", "0", "0", "1", "0"),
    (None, "OOS Products", "1", "", "", "1", "14", "5", "0", "0", "0", "1", "0"),
    (None, "Patient Progression Rate *", "12.5%", "", "0.0%", "20.0%", "4.8%", "6.2%", "3.4%", "", "0.0%", "25.0%", "0.0%"),
    (None, "AMTAGVI Infusions Performed", "4", "", "", "4", "13", "10", "7", "1", "1", "3", "0"),
    ("AMTAGVI Treatment\nTimelines", "Median Time From Enrollment Date to TTP (Days)", "37.0", "", "25.0", "42.0", "43.0", "43.0", "44.3", "", "42.0", "61.5", "35.5"),
    (None, "Median Time From TTP to AMTAGVI Infusion (Days)", "53.0", "", "", "53.0", "49.0", "47.5", "43.5", "", "52.0", "54.0", ""),
    (None, "Median Time From Final Product Delivery Date\nto AMTAGVI Infusion (Days)", "14.5", "", "", "14.5", "4.5", "5.0", "4.3", "", "10.0", "18.0", ""),
]
GROUP_FILL = [C["blue1"], C["blue2"], C["blue3"], C["blue1"]]
GROUP_ROWS = [(0, 3), (3, 6), (6, 10), (10, 13)]


def draw_table(s, x0, y0, scale=1.0, live_col=False):
    """Render the proposed scorecard. Returns the total width used."""
    wc, wm = Inches(1.30 * scale), Inches(2.62 * scale)
    wn = Inches(0.62 * scale)
    hh, hr = Inches(0.24 * scale), Inches(0.295 * scale)
    time_cols = ["Launch to\nDate", "2024", "2025", "2026\n(YTD)"]
    nat_cols = ["Top 10\nATCs *", "Top 40\nATCs *", "'New'\nATCs **"]
    qtr_cols = ["Q3'26\nQTD", "Q2'26", "Q1'26", "Q4'25"]
    live = ["Selected\nwindow"] if live_col else []
    ncols = len(time_cols) + len(live) + len(nat_cols) + len(qtr_cols)

    def cell(x, y, w, h, fill, txt, size, bold=False, col=None, ital=False,
             align=PP_ALIGN.CENTER):
        box(s, x, y, w, h, fill=fill, line=C["black"], lw=0.5)
        text(s, x + Inches(0.03), y + Inches(0.02 * scale), w - Inches(0.06), h,
             [(txt, size, bold, col or C["black"], ital)],
             align=align, anchor=MSO_ANCHOR.MIDDLE)

    # ---- block header row
    x = x0
    cell(x, y0, wc + wm, hh, C["grey_hdr"], "Uk Albert B Chandler Hospital",
         8.5 * scale, True, C["black"], True)
    x += wc + wm
    w_time = wn * (len(time_cols) + len(live))
    cell(x, y0, w_time, hh, C["grey_hdr"], "", 8 * scale)
    x += w_time
    w_nat = wn * len(nat_cols)
    cell(x, y0, w_nat, hh, C["grey_hdr"], "YTD National Metrics", 8 * scale, True,
         C["black"], True)
    x += w_nat
    cell(x, y0, wn * len(qtr_cols), hh, C["grey_hdr"], "Quarterly ATC Metrics",
         8 * scale, True, C["black"], True)

    # ---- column header row
    y = y0 + hh
    x = x0
    cell(x, y, wc, hh * 1.55, C["grey_hdr"], "Category", 8 * scale, True, C["black"], True)
    x += wc
    cell(x, y, wm, hh * 1.55, C["grey_hdr"], "Metric", 8 * scale, True, C["black"], True)
    x += wm
    for c in time_cols:
        cell(x, y, wn, hh * 1.55, C["green_hdr"], c, 7.5 * scale, True); x += wn
    for c in live:
        cell(x, y, wn, hh * 1.55, C["lime"], c, 7.5 * scale, True); x += wn
    nat_x = x
    for c in nat_cols:
        cell(x, y, wn, hh * 1.55, C["white"], c, 7.5 * scale, True); x += wn
    for c in qtr_cols:
        cell(x, y, wn, hh * 1.55, C["white"], c, 7.5 * scale, True); x += wn
    table_w = x - x0

    # ---- body
    y = y0 + hh + hh * 1.55
    for i, row in enumerate(METRICS):
        x = x0 + wc
        cell(x, y, wm, hr, C["white"], row[1], 7.2 * scale, False, C["black"],
             False, PP_ALIGN.LEFT)
        x += wm
        vals = list(row[2:6]) + ([""] if live_col else []) + list(row[6:])
        if live_col:
            vals = list(row[2:6]) + [row[2]] + list(row[6:])
        for v in vals:
            cell(x, y, wn, hr, C["white"], v, 7.5 * scale, False, C["black"],
                 False, PP_ALIGN.RIGHT)
            x += wn
        y += hr

    # ---- merged category column, drawn over the body
    for gi, (a, b) in enumerate(GROUP_ROWS):
        gy = y0 + hh + hh * 1.55 + hr * a
        gh = hr * (b - a)
        name = [m[0] for m in METRICS[a:b] if m[0]][0]
        cell(x0, gy, wc, gh, GROUP_FILL[gi], name, 8 * scale, True, C["black"],
             False, PP_ALIGN.CENTER)

    # ---- heavy box around the national block, and the red instruction
    nb = box(s, nat_x, y0 + hh, wn * len(nat_cols), hh * 1.55 + hr * len(METRICS),
             fill=None, line=C["black"], lw=2.25)
    text(s, nat_x - Inches(0.95), y0 - Inches(0.23 * scale),
         wn * len(nat_cols) + Inches(1.9), Inches(0.22),
         [("Pick one comparative arm depending on ATC", 8 * scale, True, C["red"])],
         align=PP_ALIGN.CENTER)
    return table_w, y


# ---------------------------------------------------------------- slide 4
def s4(prs):
    s = blank(prs)
    chrome(s, "P&PR Scorecard  |  The target",
           "The target: three labelled blocks, one boxed comparison arm, and colour that groups rather than decorates", 4)
    tw, ty = draw_table(s, Inches(0.62), Inches(2.15), scale=0.96)
    text(s, Inches(0.62), ty + Inches(0.12), Inches(11.9), Inches(0.8),
         [[("* Patient Progression Rate = patient related drop-offs after manufacturing "
            "start, divided by manufacturing starts.  * Top 10 and Top 40 ATCs are the "
            "highest enrolling centres during the specific timeframe.", 8.5, False, C["ink"])],
          [("** 'New' refers to ATCs authorized and onboarded in the 2025 calendar year.  "
            "* TTP cancellations are estimated until the Infinity snapshot history is "
            "connected.", 8.5, False, C["ink"])]], space=2)


# ---------------------------------------------------------------- slide 5
def s5(prs):
    s = blank(prs)
    chrome(s, "P&PR Scorecard  |  What changes",
           "Four changes take the automated table from a data dump back to a document", 5)
    rows = [("Two-row header with named blocks",
             "Thirteen equal columns with no grouping",
             "The centre's own figures, the national comparison and the quarterly trend "
             "are three different questions. Labelling the blocks tells the eye which one "
             "it is reading."),
            ("Category column shaded by group",
             "A white column of repeated words",
             "Four coloured blocks let someone find the funnel stage they care about "
             "without reading thirteen row labels."),
            ("One boxed comparison arm",
             "Three benchmark columns shown side by side",
             "Kolin's own note on the template says pick one arm depending on the ATC. "
             "A large centre compared against 'New' ATCs is a misleading number."),
            ("Diagnostics off the printed view",
             "Undated and After as-of shown as columns",
             "Those two exist so the pipeline can prove no event was lost. They are "
             "build checks, not review metrics, and they do not belong in front of a doctor.")]
    y = Inches(2.02)
    for a, b, why in rows:
        box(s, Inches(0.62), y, Inches(3.35), Inches(1.06), fill=C["panel"],
            line=C["olive"], lw=0.75)
        text(s, Inches(0.78), y + Inches(0.13), Inches(3.05), Inches(0.85),
             [(a, 11.5, True, C["olive"])], anchor=MSO_ANCHOR.MIDDLE)
        box(s, Inches(4.12), y, Inches(3.05), Inches(1.06), fill=C["white"],
            line=C["faint"], lw=0.5)
        text(s, Inches(4.27), y + Inches(0.13), Inches(2.78), Inches(0.85),
             [(b, 10.5, False, C["faint"])], anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(7.44), y + Inches(0.10), Inches(5.1), Inches(0.95),
             [(why, 10.5, False, C["ink"])], anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.17)
    text(s, Inches(0.62), Inches(1.72), Inches(3.35), Inches(0.25),
         [("PROPOSED", 9, True, C["navy"])])
    text(s, Inches(4.12), Inches(1.72), Inches(3.05), Inches(0.25),
         [("CURRENT BUILD", 9, True, C["faint"])])
    text(s, Inches(7.44), Inches(1.72), Inches(5.1), Inches(0.25),
         [("WHY IT MATTERS", 9, True, C["navy"])])


# ---------------------------------------------------------------- slide 6
def s6(prs):
    s = blank(prs)
    chrome(s, "P&PR Scorecard  |  The one addition",
           "One added column, the date window, answers the question the fixed columns cannot", 6)
    tw, ty = draw_table(s, Inches(0.62), Inches(1.98), scale=0.79, live_col=True)
    box(s, Inches(0.62), Inches(5.70), Inches(11.9), Inches(1.10),
        fill=C["white"], line=C["navy"], lw=1.25)
    text(s, Inches(0.95), Inches(5.88), Inches(11.3), Inches(0.88),
         [[("Why one extra column earns its place", 12, True, C["navy"])],
          [("The fixed columns cannot answer \"what has this centre done since we changed "
            "the protocol in March\". Two date boxes drive the lime column only, so 2024 "
            "and 2025 never move. A column headed 2024 that reads zero because of a "
            "control somewhere else on the page is how people stop trusting every other "
            "number on it.", 11, False, C["ink"])]], space=5)


# ---------------------------------------------------------------- slide 7
def s7(prs):
    s = blank(prs)
    chrome(s, "P&PR Scorecard  |  Next",
           "Three things to settle before this goes in front of a treatment centre", 7)
    items = [("Confirm the comparison arm rule",
              "The template says pick one arm depending on the ATC. We need the rule "
              "written down: which centres see Top 10, which see Top 40, which see New. "
              "Today the dashboard shows all three."),
             ("Replace the estimated cancellation metric",
              "TTPs Cancelled is still a stand-in. The real logic reads Infinity's "
              "snapshot history and is written and waiting on one export. Until then the "
              "asterisk and the footnote stay."),
             ("Shadow-run three centres Kolin knows",
              "Reproduce three decks he has already made by hand and walk the "
              "differences with him before any centre sees this. The first disputed cell "
              "decides whether the other twelve are believed.")]
    y = Inches(2.15)
    for i, (h, b) in enumerate(items, 1):
        box(s, Inches(0.62), y, Inches(0.42), Inches(0.42), fill=C["olive"])
        text(s, Inches(0.62), y + Inches(0.07), Inches(0.42), Inches(0.3),
             [(str(i), 14, True, C["white"])], align=PP_ALIGN.CENTER)
        text(s, Inches(1.28), y + Inches(0.02), Inches(11.2), Inches(0.32),
             [(h, 15, True, C["navy"])])
        text(s, Inches(1.28), y + Inches(0.44), Inches(10.9), Inches(0.8),
             [(b, 11.5, False, C["ink"])])
        y += Inches(1.42)
    text(s, Inches(0.62), Inches(6.55), Inches(11.9), Inches(0.4),
         [("Every figure shown in this deck is from the current build on real Infinity "
           "data for Uk Albert B Chandler Hospital.", 10, False, C["faint"])])


def main():
    prs = deck()
    for fn in (s1, s2, s3, s4, s5, s6, s7):
        fn(prs)
    out = "P&PR Scorecard - table design proposal.pptx"
    prs.save(out)
    print("wrote", out, "|", len(prs.slides.__iter__.__self__._sldIdLst), "slides")


if __name__ == "__main__":
    main()
